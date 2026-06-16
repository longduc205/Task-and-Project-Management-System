#!/usr/bin/env python3
"""
build_pptx.py — Build Group 14's INS3044 PPTX deck
Run: python3 docs/slides/build_pptx.py

Design system (Cartesian):
  --bg:         #F5EFE6  (warm parchment)
  --bg-elev:    #EDE4D3  (card fill)
  --ink:        #2B2419  (deep walnut — primary text)
  --ink-soft:   #5C4F3F  (secondary text)
  --muted:      #8A7A66  (captions)
  --rule:       #D9CDB8  (hairlines)
  --accent:     #8C5A2E  (burnt sienna — use sparingly)
  --accent-soft:#E8D5BC  (soft accent fill)
  --ok:         #5B6B47  (olive)
  --warn:       #A65D3F  (rust)

Fonts:
  Title: 'Playfair Display', Georgia (fallback)
  Body:  Calibri (universal)
  Mono:  Consolas (tabular numbers)
"""

import os, sys, math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colour palette ────────────────────────────────────────────────────────────
BG         = RGBColor(0xF5, 0xEF, 0xE6)
BG_ELEV    = RGBColor(0xED, 0xE4, 0xD3)
INK        = RGBColor(0x2B, 0x24, 0x19)
INK_SOFT   = RGBColor(0x5C, 0x4F, 0x3F)
MUTED      = RGBColor(0x8A, 0x7A, 0x66)
RULE       = RGBColor(0xD9, 0xCD, 0xB8)
ACCENT     = RGBColor(0x8C, 0x5A, 0x2E)
ACCENT_SFT = RGBColor(0xE8, 0xD5, 0xBC)
OK         = RGBColor(0x5B, 0x6B, 0x47)
WARN       = RGBColor(0xA6, 0x5D, 0x3F)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

# ── Slide dimensions (16:9) ──────────────────────────────────────────────────
SW = Inches(13.333)
SH = Inches(7.5)

# ── Helpers ──────────────────────────────────────────────────────────────────
def rgb_hex(rgb: RGBColor):
    return rgb  # keep as RGBColor

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    """Add a filled rectangle shape."""
    shp = slide.shapes.add_shape(
        1, SW * l if l < 1 else l,
           SH * t if t < 1 else t,
           SW * w if w < 1 else w,
           SH * h if h < 1 else h)
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w or 0.75)
    else:
        shp.line.fill.background()
    return shp

def add_text(slide, text, l, t, w, h=None,
             font="Calibri", size=14, bold=False, italic=False,
             color=INK, align=PP_ALIGN.LEFT, wrap=True,
             v_anchor=None, fill=None, line=None, line_w=None):
    """Add a text box."""
    txb = slide.shapes.add_textbox(
        SW * l if l < 1 else l,
        SH * t if t < 1 else t,
        SW * w if w < 1 else w,
        SH * (h or 0.4) if h and h < 1 else (h or Inches(0.4)))
    if fill:
        add_rect(slide, l, t, w, h or 0.4, fill=fill, line=line, line_w=line_w)
    tf = txb.text_frame
    tf.word_wrap = wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_para(tf, text, font="Calibri", size=14, bold=False, italic=False,
             color=INK, align=PP_ALIGN.LEFT, space_before=0, space_after=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def add_table(slide, data, l, t, w, h,
              header_fill=ACCENT, header_text=WHITE,
              row_fills=None,  # list of RGBColor or None
              font="Calibri", font_size=11,
              col_widths=None, align=PP_ALIGN.LEFT):
    """Add a styled table. data[0] = header row."""
    rows = len(data)
    cols = len(data[0])
    tbl = slide.shapes.add_table(rows, cols,
        SW*l if l < SW else l,
        SH*t if t < SH else t,
        SW*w if w < SW else w,
        SH*h if h < SH else h).table

    # Column widths
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(cw / total * (SW*w if w < SW else w))

    for r, row in enumerate(data):
        for c, cell_text in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = str(cell_text)
            run.font.name = font
            run.font.size = Pt(font_size)
            run.font.color.rgb = WHITE if r == 0 else INK
            # Fill
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            elif row_fills and r-1 < len(row_fills):
                if row_fills[r-1]:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = row_fills[r-1]
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG
    return tbl

def section_divider(prs, section_num, roman, title, cline=""):
    """Full-bleed section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Background
    shp = slide.shapes.add_shape(1, 0, 0, SW, SH)
    shp.fill.solid()
    shp.fill.fore_color.rgb = BG
    shp.line.fill.background()
    # Hairline top
    add_rect(slide, 0, 0, 1, 0.008, line=RULE, line_w=1)
    # Hairline bottom
    add_rect(slide, 0, 0.992, 1, 0.008, line=RULE, line_w=1)
    # § mark
    add_text(slide, "§", 0.8, 1.5, 2, 1.5,
             font="Georgia", size=120, italic=True, color=ACCENT,
             align=PP_ALIGN.LEFT)
    # Roman numeral
    add_text(slide, roman, 2.8, 1.7, 2, 1,
             font="Georgia", size=48, italic=True, color=ACCENT_SFT,
             align=PP_ALIGN.LEFT)
    # Section title
    add_text(slide, title, 0.8, 3.4, 11, 1,
             font="Georgia", size=44, bold=True, color=INK,
             align=PP_ALIGN.LEFT)
    # CLO note
    if cline:
        add_text(slide, cline, 0.8, 4.5, 10, 0.5,
                 font="Calibri", size=14, italic=True, color=MUTED,
                 align=PP_ALIGN.LEFT)
    # Footer rule + text
    add_rect(slide, 0.8, 7.0, 11.733, 0.003, line=RULE, line_w=0.75)
    add_text(slide, "Group 14 · INS3044 · June 2026",
             0.8, 7.05, 5, 0.3,
             font="Calibri", size=10, color=MUTED)
    return slide

def content_slide(prs, title, subtitle=None, eyebrow=None,
                  notes="", accent_one=False, slide_num=None, total=None):
    """Standard content slide with header, footer."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Background
    shp = slide.shapes.add_shape(1, 0, 0, SW, SH)
    shp.fill.solid()
    shp.fill.fore_color.rgb = BG
    shp.line.fill.background()
    # Top hairline
    add_rect(slide, 0, 0, 1, 0.005, line=RULE, line_w=0.75)
    # Eyebrow
    if eyebrow:
        add_text(slide, eyebrow.upper(), 0.7, 0.15, 8, 0.35,
                 font="Calibri", size=10, color=MUTED)
    # Title
    t_color = ACCENT if accent_one else INK
    add_text(slide, title, 0.7, 0.45, 12, 0.8,
             font="Georgia", size=30, bold=True, color=t_color,
             align=PP_ALIGN.LEFT)
    # Subtitle
    if subtitle:
        add_text(slide, subtitle, 0.7, 1.05, 12, 0.4,
                 font="Calibri", size=14, italic=True, color=INK_SOFT)
    # Bottom footer rule
    add_rect(slide, 0.7, 7.2, 11.933, 0.003, line=RULE, line_w=0.75)
    # Footer left
    add_text(slide, "Group 14 · INS3044 · June 2026",
             0.7, 7.25, 5, 0.3,
             font="Calibri", size=10, color=MUTED)
    # Footer right slide number
    if slide_num and total:
        add_text(slide, f"{slide_num} / {total}",
                 12.3, 7.25, 1, 0.3,
                 font="Calibri", size=10, color=MUTED, align=PP_ALIGN.RIGHT)
    # Notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    return slide

def kpi_card(slide, l, t, w, h, label, value, sub="", accent=False):
    """KPI card: rounded rect outline, big number, small label."""
    fill_c = ACCENT_SFT if accent else BG_ELEV
    add_rect(slide, l, t, w, h, fill=fill_c,
             line=ACCENT if accent else RULE, line_w=0.75)
    add_text(slide, value, l+0.1, t+0.05, w-0.2, h*0.6,
             font="Georgia", size=28, bold=True,
             color=ACCENT if accent else INK, align=PP_ALIGN.CENTER)
    add_text(slide, label, l+0.05, t+h*0.55, w-0.1, h*0.25,
             font="Calibri", size=10, color=INK_SOFT, align=PP_ALIGN.CENTER)
    if sub:
        add_text(slide, sub, l+0.05, t+h*0.75, w-0.1, h*0.2,
                 font="Calibri", size=9, italic=True, color=MUTED,
                 align=PP_ALIGN.CENTER)

# ── Build deck ────────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    TOTAL = 24  # estimated; tracked as we go
    slide_num = [0]

    def next_num():
        slide_num[0] += 1
        return slide_num[0]

    notes_master = ""

    # ════════════════════════════════════════════════════════════════
    # SLIDE 1 — COVER
    # ════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Background
    shp = s.shapes.add_shape(1, 0, 0, SW, SH)
    shp.fill.solid(); shp.fill.fore_color.rgb = BG; shp.line.fill.background()
    # Top hairline
    add_rect(s, 0, 0, 1, 0.005, line=RULE, line_w=1)
    # Course kicker
    add_text(s, "INS3044 · IT Project Management", 0.9, 0.6, 11, 0.45,
             font="Calibri", size=14, color=MUTED)
    # Main title
    add_text(s, "Task & Project", 0.9, 1.2, 11.5, 1.1,
             font="Georgia", size=62, bold=True, color=INK)
    add_text(s, "Management System", 0.9, 2.1, 11.5, 0.9,
             font="Georgia", size=62, bold=True, color=INK)
    # Tagline
    add_text(s, "An unhurried 14-week build",
             0.9, 3.1, 9, 0.5,
             font="Georgia", size=22, italic=True, color=INK_SOFT)
    # Hairline rule
    add_rect(s, 0.9, 3.65, 11.533, 0.003, line=RULE, line_w=1)
    # Member cards
    members = [
        ("PM · Leader",        "Nguyễn Long Đức",  "23070435"),
        ("Developer",          "Phạm Hồ Bảo",       "23070455"),
        ("UI / UX Designer",   "Kiều Bá Thịnh",      "23070247"),
        ("Quality Assurance",  "Đỗ Huy Hiếu",        "23070325"),
    ]
    for i, (role, name, sid) in enumerate(members):
        x = 0.9 + i * 2.9
        add_rect(s, x, 3.85, 2.7, 1.1, fill=BG_ELEV, line=RULE, line_w=0.75)
        add_text(s, role,  x+0.1, 3.9,  2.5, 0.3, font="Calibri", size=9,  color=MUTED)
        add_text(s, name,  x+0.1, 4.15, 2.5, 0.35, font="Georgia", size=14, bold=True, color=INK)
        add_text(s, sid,   x+0.1, 4.55, 2.5, 0.3, font="Calibri", size=10, color=MUTED)
    # Footer
    add_rect(s, 0.9, 6.6, 11.533, 0.003, line=RULE, line_w=0.75)
    add_text(s, "Instructors: Dr Nguyễn Phương Anh · MS Đỗ Tiến Thành",
             0.9, 6.65, 6, 0.3, font="Calibri", size=10, color=MUTED)
    add_text(s, "Faculty of Applied Sciences · VNU · June 2026",
             0.9, 6.9, 6, 0.3, font="Calibri", size=10, color=MUTED)
    s.notes_slide.notes_text_frame.text = (
        "Welcome everyone. Group 14 consists of four members: Đức as project manager, "
        "Bảo on development, Thịnh handling UI and UX, and Hiếu leading QA. We've spent "
        "14 weeks building a web-based Task and Project Management System from scratch, "
        "applying every major discipline of IT project management along the way. The tagline "
        "'an unhurried 14-week build' captures our approach: we prioritised sound process "
        "over speed, which is exactly what INS3044 asks us to demonstrate. The talk is "
        "structured across eight sections — WBS, PERT, CPM, budgeting, risk, communication, "
        "EVM monitoring, and the prototype demo. Questions are welcome throughout."
    )

    # ════════════════════════════════════════════════════════════════
    # SLIDE 2 — AGENDA
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = content_slide(prs,
        title="Agenda",
        subtitle="Eight sections, one project",
        eyebrow="Overview",
        notes="Our presentation is organised into eight sections, mirroring the structure of a "
              "full project management cycle. We'll open with scope and objectives, move through "
              "the planning artefacts — WBS, schedule, budget — then the execution controls "
              "of risk and communication, finish with EVM monitoring and the live demo. Each "
              "section maps to specific CLOs in the INS3044 syllabus.",
        slide_num=sn, total=TOTAL)
    sections = [
        ("§ I",  "Scope & Objectives",       "CLO 1, 2"),
        ("§ II", "Work Breakdown Structure", "CLO 2, 5"),
        ("§ III","Timeline & Milestones",    "CLO 2, 4, 5"),
        ("§ IV", "Resource Management",      "CLO 2, 8, 9"),
        ("§ V",  "Risk Management",          "CLO 3, 6"),
        ("§ VI", "Communication",            "CLO 4, 5, 8, 9, 12"),
        ("§ VII","Monitoring & EVM",         "CLO 7"),
        ("§ VIII","Prototype & Lessons",      "CLO 7, 8, 9"),
    ]
    col_w = 5.5
    for i, (num, title, clo) in enumerate(sections):
        col = i % 2
        row = i // 2
        x = 0.7 + col * col_w
        y = 1.65 + row * 1.15
        add_rect(s, x, y, 5.2, 0.95, fill=BG_ELEV, line=RULE, line_w=0.75)
        add_text(s, num, x+0.15, y+0.05, 0.7, 0.85,
                 font="Georgia", size=28, italic=True, color=ACCENT)
        add_text(s, title, x+0.9, y+0.1, 4.1, 0.4,
                 font="Georgia", size=16, bold=True, color=INK)
        add_text(s, clo, x+0.9, y+0.55, 4.1, 0.3,
                 font="Calibri", size=11, color=MUTED)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 3 — § I DIVIDER
    # ════════════════════════════════════════════════════════════════
    section_divider(prs, 3, "I", "Scope & Objectives",
                    "CLO 1, 2 · Ch. 1, 2, 3")
    prs.slides[-1].notes_slide.notes_text_frame.text = (
        "Section I sets the foundation. We'll cover the business context, define what's "
        "in scope and what's explicitly out of scope, and walk through our five SMART "
        "objectives that tie directly to the INS3044 CLOs."
    )

    # ════════════════════════════════════════════════════════════════
    # SLIDE 4 — PROJECT KPIs
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = content_slide(prs,
        title="Project at a Glance",
        eyebrow="Scope & Objectives",
        notes="These four KPIs define the project baseline. Budget at Completion is 33.6 million "
              "VND, covering labour, overhead, G&A, contingency, and management reserve. The "
              "critical path is 43 working days — we'll see how we derived that in Section III. "
              "The risk EMV of 3.08 million VND is fully covered by contingency plus a "
              "management reserve top-up. And we have five prototype features to deliver.",
        slide_num=sn, total=TOTAL)
    kpis = [
        ("33,600,000 VND",  "Budget at Completion (BAC)",            True),
        ("43 working days", "Critical Path Duration",                  False),
        ("3,080,000 VND",  "Total Risk EMV (Cost)",                  False),
        ("5 core features", "Prototype Feature Areas",                False),
    ]
    for i, (val, lbl, acc) in enumerate(kpis):
        x = 0.7 + i * 3.1
        kpi_card(s, x, 1.7, 2.9, 1.8, lbl, val, accent=acc)
    # Sub-kpis
    add_text(s, "BAC breakdown: Labour 24M · OH 20% · G&A reallocated · Contingency 2.4M · Mgmt reserve 0.68M",
             0.7, 3.7, 12, 0.3, font="Calibri", size=11, italic=True, color=MUTED)
    add_text(s, "Risk EMV coverage: Contingency 2,400,000 + Management reserve 680,000 = 3,080,000 VND ✓",
             0.7, 4.1, 12, 0.3, font="Calibri", size=11, italic=True, color=MUTED)
    # Team summary
    add_rect(s, 0.7, 4.6, 12, 0.003, line=RULE, line_w=0.5)
    add_text(s, "Team: Nguyễn Long Đức (23070435, PM) · Phạm Hồ Bảo (23070455, Dev) · "
                "Kiều Bá Thịnh (23070247, UI/UX) · Đỗ Huy Hiếu (23070325, QA)",
             0.7, 4.7, 12, 0.4, font="Calibri", size=11, color=INK_SOFT)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 5 — IN-SCOPE / OUT-OF-SCOPE
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = content_slide(prs,
        title="Scope Boundaries",
        eyebrow="Scope & Objectives",
        notes="The left column shows what's explicitly in scope — agreed by all four team members "
              "and traceable to the project registration. The right column documents what we "
              "deliberately excluded and why. This scope lock is our first line of defence "
              "against scope creep, which is Risk R2 in our risk register.",
        slide_num=sn, total=TOTAL)
    # In-scope
    add_rect(s, 0.7, 1.55, 5.8, 0.4, fill=OK, line=None)
    add_text(s, "IN SCOPE", 0.7, 1.55, 5.8, 0.4,
             font="Calibri", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    in_scope = [
        "Project CRUD (create, read, update, delete)",
        "Task CRUD with user assignment",
        "Status workflow: To Do → In Progress → Done",
        "Task deadline with overdue highlighting",
        "Task progress % (0–100%) per task",
        "Project dashboard (counts & progress)",
        "Role-based UI: Admin / PM / Member",
        "Responsive single-page web application",
        "Functional prototype (≥ 3 features live)",
        "Full PM documentation (WBS, CPM, budget, risk, EVM)",
    ]
    for i, item in enumerate(in_scope):
        add_text(s, "·  " + item, 0.7, 2.0+i*0.41, 5.8, 0.38,
                 font="Calibri", size=10, color=INK)
    # Out-of-scope
    add_rect(s, 6.8, 1.55, 5.8, 0.4, fill=WARN, line=None)
    add_text(s, "OUT OF SCOPE", 6.8, 1.55, 5.8, 0.4,
             font="Calibri", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    out_scope = [
        "Native iOS / Android application",
        "Real-time chat or messaging",
        "Billing, invoicing, payment processing",
        "Third-party integrations (Slack, GitHub, Outlook)",
        "Multi-tenancy (multiple organisations)",
        "Post-launch maintenance or support",
        "Automated email / push notifications",
    ]
    for i, item in enumerate(out_scope):
        add_text(s, "✗  " + item, 6.8, 2.0+i*0.41, 5.8, 0.38,
                 font="Calibri", size=10, color=INK)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 6 — SMART OBJECTIVES
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = content_slide(prs,
        title="SMART Objectives",
        eyebrow="Scope & Objectives",
        notes="Each of our five objectives is Specific, Measurable, Achievable, Relevant, and "
              "Time-bound — the SMART framework. OBJ-1 covers the prototype deliverables. "
              "OBJ-2 covers our 480-hour budget. OBJ-3 covers quality at UAT. OBJ-4 covers "
              "risk management and the EMV-driven contingency. OBJ-5 covers the final submission. "
              "Every objective is tagged to one or more CLOs, so we can demonstrate coverage.",
        slide_num=sn, total=TOTAL)
    headers = ["#", "Objective", "CLO"]
    data = [
        headers,
        ["OBJ-1", "Deliver working prototype: task CRUD, assignment, deadline tracking,\n"
                  "progress monitoring for 3 roles — demo at W10, submit at W14", "CLO 1, 2"],
        ["OBJ-2", "Complete all 16 WBS Level 3 packages within 480 person-hours (±10%)\n"
                  "by Week 14; tracked weekly in Trello", "CLO 2, 4"],
        ["OBJ-3", "Pass UAT with ≥ 80% test coverage and ≤ 10 defects/KLOC;\n"
                  "UAT sign-off at Week 12", "CLO 3, 6"],
        ["OBJ-4", "Identify ≥ 6 risks; maintain contingency reserve ≥ Total Risk EMV;\n"
                  "risk register updated bi-weekly", "CLO 3, 6"],
        ["OBJ-5", "Submit full ZIP package (report, prototype, PPTX+PDF, evidence)\n"
                  "before Week 14 deadline", "CLO 7, 8, 9"],
    ]
    col_w = [0.65, 10.55, 0.9]
    add_table(s, data, 0.7, 1.55, 12.3, 3.8,
              col_widths=col_w, font_size=11)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 7 — § II DIVIDER
    # ════════════════════════════════════════════════════════════════
    section_divider(prs, 7, "II", "Work Breakdown Structure",
                    "CLO 2, 5 · Ch. 4, 5")
    prs.slides[-1].notes_slide.notes_text_frame.text = (
        "Section II covers the Work Breakdown Structure. The WBS is our definitive scope "
        "baseline — every deliverable maps to exactly one Level 3 work package. We'll look "
        "at the three-level hierarchy and the WBS dictionary."
    )

    # ════════════════════════════════════════════════════════════════
    # SLIDE 8 — WBS TREE
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = content_slide(prs,
        title="WBS — Three-Level Hierarchy",
        eyebrow="Work Breakdown Structure",
        notes="The WBS decomposes the project into manageable chunks. Level 1 is the project "
              "itself. Level 2 has seven branches — Project Management, Requirements, Design, "
              "Development, QA, Documentation, and Presentation. Level 3 gives us 16 work "
              "packages that together account for 100% of the scope — verified against the "
              "scope statement. Total estimated effort is 480 person-hours.",
        slide_num=sn, total=TOTAL)

    wbs_items = [
        ("1.0",  "Task & Project Management System", True,  False),
        ("1.1",  "Project Management",                 False, True),
        ("1.1.1","  Project Charter & Kickoff",         False, False),
        ("1.1.2","  Status Meetings (12 × 2 h)",        False, False),
        ("1.1.3","  Risk Reviews",                       False, False),
        ("1.1.4","  Final Report Writing",              False, False),
        ("1.2",  "Requirements & Analysis",             False, True),
        ("1.2.1","  Stakeholder Interviews",            False, False),
        ("1.2.2","  Requirements Document",              False, False),
        ("1.2.3","  Requirements Sign-off",              False, False),
        ("1.3",  "System Design",                        False, True),
        ("1.3.1","  UI/UX Design (Figma)",               False, False),
        ("1.3.2","  Database / Data-Model Schema",       False, False),
        ("1.3.3","  API / Service Contract",             False, False),
        ("1.4",  "Development",                          False, True),
        ("1.4.1","  Frontend (React + Vite + TS)",       False, False),
        ("1.4.2","  Backend / Persistence Layer",        False, False),
        ("1.4.3","  Integration & Wiring",               False, False),
        ("1.5",  "Testing & QA",                          False, True),
        ("1.5.1","  Unit Tests (Vitest)",                 False, False),
        ("1.5.2","  Integration Tests",                  False, False),
        ("1.5.3","  User Acceptance Test",               False, False),
        ("1.6",  "Documentation",                        False, True),
        ("1.6.1","  User Guide",                         False, False),
        ("1.6.2","  Technical Documentation",            False, False),
        ("1.7",  "Presentation & Submission",            False, True),
        ("1.7.1","  Slide Deck",                         False, False),
        ("1.7.2","  Rehearsal",                          False, False),
    ]

    for i, (wbs_id, name, is_root, is_l2) in enumerate(wbs_items):
        row = i
        x_indent = 0.7
        if is_root:
            x_indent = 0.7
        elif is_l2:
            x_indent = 0.9
        else:
            x_indent = 1.5
        fsize = 12 if is_root else (11 if is_l2 else 10)
        fbold = is_root or is_l2
        fcol = INK if is_root else (INK if is_l2 else INK_SOFT)
        add_rect(s, x_indent, 1.55+row*0.215, 4.5, 0.2,
                 fill=BG_ELEV if is_l2 else None,
                 line=RULE if is_l2 else None, line_w=0.5)
        add_text(s, wbs_id, x_indent+0.05, 1.55+row*0.215, 0.55, 0.2,
                 font="Calibri", size=fsize-1, bold=True, color=ACCENT)
        add_text(s, name, x_indent+0.6, 1.55+row*0.215, 3.9, 0.2,
                 font="Calibri", size=fsize, bold=fbold, color=fcol)

    # Hours summary
    add_rect(s, 6.5, 1.55, 6.3, 0.003, line=RULE, line_w=0.5)
    add_text(s, "WBS Level 3 Hours Summary", 6.5, 1.65, 6.3, 0.4,
             font="Georgia", size=14, bold=True, color=INK)
    hrs = [("1.1 PM", "60 h"), ("1.2 Req", "50 h"), ("1.3 Design", "80 h"),
           ("1.4 Dev", "180 h"), ("1.5 QA", "60 h"), ("1.6 Docs", "30 h"),
           ("1.7 Pres", "20 h")]
    for i, (lbl, val) in enumerate(hrs):
        row = i // 2; col = i % 2
        x = 6.5 + col * 3.1; y = 2.2 + row * 0.55
        add_rect(s, x, y, 2.9, 0.45, fill=BG_ELEV, line=RULE, line_w=0.5)
        add_text(s, lbl, x+0.1, y+0.03, 1.6, 0.38, font="Calibri", size=10, color=MUTED)
        add_text(s, val, x+1.7, y+0.03, 1.1, 0.38, font="Georgia", size=14, bold=True, color=INK)
    add_rect(s, 6.5, 4.2, 6.3, 0.003, line=RULE, line_w=0.75)
    add_text(s, "Total Labor: 480 hours × 50,000 VND/h = 24,000,000 VND",
             6.5, 4.3, 6.3, 0.4, font="Calibri", size=11, bold=True, color=INK)
    add_text(s, "→ See Budget Rollup in Section IV",
             6.5, 4.75, 6.3, 0.3, font="Calibri", size=11, italic=True, color=MUTED)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 9 — § III DIVIDER
    # ════════════════════════════════════════════════════════════════
    section_divider(prs, 9, "III", "Timeline & Milestones",
                    "CLO 2, 4, 5 · Ch. 4, 8, 10")
    prs.slides[-1].notes_slide.notes_text_frame.text = (
        "Section III is the schedule heart of the project. We'll cover PERT three-point "
        "estimation for six activities, the full CPM forward and backward pass, and the "
        "critical path identification. The Gantt overview shows all 16 work packages "
        "mapped against the project timeline with six key milestones."
    )

    # ════════════════════════════════════════════════════════════════
    # SLIDE 10 — GANTT OVERVIEW (built with shapes)
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = content_slide(prs,
        title="Gantt Chart — 14-Week Timeline",
        eyebrow="Timeline & Milestones",
        notes="The Gantt chart shows all 16 Level 3 work packages across the 14-week project. "
              "The critical path activities are highlighted in burnt sienna — they have zero "
              "float and any delay on them directly delays the project end. The non-critical "
              "activities are in muted brown. Six diamond milestones mark key delivery points. "
              "Activity F (Backend Development) has 8 days of float — it can slip without "
              "affecting the critical path.",
        slide_num=sn, total=TOTAL)

    # Header row — weeks
    add_text(s, "WBS", 0.7, 1.55, 1.6, 0.35,
             font="Calibri", size=9, bold=True, color=MUTED)
    for w in range(1, 15):
        x = 2.35 + (w-1) * 0.76
        add_rect(s, x, 1.55, 0.73, 0.35, fill=BG_ELEV, line=RULE, line_w=0.5)
        add_text(s, f"W{w}", x, 1.55, 0.73, 0.35,
                 font="Calibri", size=9, bold=True, color=INK, align=PP_ALIGN.CENTER)

    # WBS rows: (wbs_id, name, start_w, duration, is_critical, owner)
    gantt_rows = [
        ("1.1.1", "Charter & Kickoff",     1, 2,  True,  "Đức"),
        ("1.1.2", "Status Meetings",         1, 14, False, "Đức"),
        ("1.1.3", "Risk Reviews",           1, 14, False, "Đức"),
        ("1.2.1", "Stakeholder Interviews", 2, 1,  True,  "Đức"),
        ("1.2.2", "Requirements Doc",       2, 2,  True,  "Bảo"),
        ("1.2.3", "Req Sign-off",           3, 1,  True,  "Đức"),
        ("1.3.1", "UI/UX Design",           3, 3,  True,  "Thịnh"),
        ("1.3.2", "DB Schema",               3, 2,  False, "Bảo"),
        ("1.3.3", "API Contract",            4, 2,  False, "Bảo"),
        ("1.4.1", "Frontend Dev",            5, 5,  True,  "Bảo+Thịnh"),
        ("1.4.2", "Backend Dev",             5, 4,  False, "Bảo"),
        ("1.4.3", "Integration",             8, 2,  True,  "Bảo"),
        ("1.5.1", "Unit Tests",              7, 4,  False, "Hiếu"),
        ("1.5.2", "Integration Tests",       9, 3,  False, "Hiếu"),
        ("1.5.3", "UAT",                    12, 1,  False, "Hiếu"),
        ("1.6.1", "User Guide",             11, 3,  False, "Thịnh"),
        ("1.6.2", "Tech Docs",              11, 3,  False, "Bảo"),
        ("1.7.1", "Slide Deck",             12, 2,  False, "Thịnh"),
        ("1.7.2", "Rehearsal",              13, 2,  False, "Đức"),
    ]
    # Milestones: (week, label)
    milestones = [
        (1,  "M1: Charter\nApproved"),
        (3,  "M2: Req\nSign-off"),
        (5,  "M3: Design\nReview"),
        (10, "M4: Proto\nDemo"),
        (12, "M5: Testing\nComplete"),
        (14, "M6: Final\nSubmission"),
    ]

    for i, (wid, name, start, dur, crit, owner) in enumerate(gantt_rows):
        row_y = 1.95 + i * 0.255
        # WBS label
        add_text(s, wid, 0.7, row_y, 0.55, 0.24,
                 font="Calibri", size=7, bold=True, color=ACCENT)
        add_text(s, name[:18], 1.28, row_y, 1.05, 0.24,
                 font="Calibri", size=7, color=INK)
        # Bar
        bar_x = 2.35 + (start-1) * 0.76
        bar_w = dur * 0.76 - 0.04
        bar_fill = ACCENT if crit else RGBColor(0xC4, 0xB4, 0x9A)
        add_rect(s, bar_x, row_y+0.04, bar_w, 0.18,
                 fill=bar_fill, line=None)
        add_text(s, owner, bar_x+0.02, row_y+0.04, bar_w-0.04, 0.18,
                 font="Calibri", size=6, color=WHITE if crit else INK,
                 align=PP_ALIGN.CENTER)

    # Milestone diamonds on bottom axis
    axis_y = 1.95 + len(gantt_rows) * 0.255 + 0.05
    add_rect(s, 2.35, axis_y, 10.64, 0.003, line=RULE, line_w=0.5)
    for w, lbl in milestones:
        mx = 2.35 + (w-1)*0.76 + 0.36
        # Diamond shape (rotate a square)
        d = s.shapes.add_shape(12, mx-0.1, axis_y+0.02, 0.2, 0.2)  # star
        d.fill.solid(); d.fill.fore_color.rgb = ACCENT
        d.line.fill.background()
        add_text(s, lbl, mx-0.4, axis_y+0.22, 0.85, 0.5,
                 font="Calibri", size=7, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Legend
    add_rect(s, 0.7, 7.05, 0.3, 0.18, fill=ACCENT, line=None)
    add_text(s, "Critical path", 1.05, 7.05, 1.8, 0.2, font="Calibri", size=9, color=INK)
    add_rect(s, 3.0, 7.05, 0.3, 0.18, fill=RGBColor(0xC4,0xB4,0x9A), line=None)
    add_text(s, "Non-critical", 3.35, 7.05, 1.5, 0.2, font="Calibri", size=9, color=INK)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 11 — PERT THREE-POINT
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = content_slide(prs,
        title="PERT Three-Point Estimates",
        eyebrow="Timeline & Milestones",
        notes="PERT uses three-point estimation to account for uncertainty. For each activity "
              "we record an Optimistic (O), Most Likely (M), and Pessimistic (P) duration. "
              "The weighted expected time TE = (O + 4M + P) / 6 gives more weight to the "
              "most likely estimate. Standard deviation σ = (P − O) / 6 tells us how certain "
              "we are. Activity B — UI/UX Design — is the work-through example shown here.",
        slide_num=sn, total=TOTAL)
    pert_data = [
        ["Activity", "Description",        "O (d)", "M (d)", "P (d)", "TE = (O+4M+P)/6", "σ = (P−O)/6", "Notes"],
        ["A",        "Requirements",       "3",      "5",      "9",      "5.33 d",           "1.00 d",      "Critical path"],
        ["B",        "UI/UX Design",       "5",      "10",     "18",     "10.50 d",          "2.17 d",      "Critical path"],
        ["E",        "Frontend Dev",       "10",     "15",     "22",     "15.33 d",          "2.00 d",      "Critical path"],
        ["G",        "Integration",         "3",      "5",      "8",      "5.17 d",           "0.83 d",      "Critical path"],
        ["H",        "Testing",             "5",      "8",      "14",     "8.50 d",           "1.50 d",      "Critical path"],
        ["F",        "Backend Dev",         "8",      "12",     "20",     "12.67 d",          "2.00 d",      "Non-critical; 8 d float"],
    ]
    add_table(s, pert_data, 0.7, 1.55, 12.3, 2.6,
              col_widths=[0.7, 1.7, 0.65, 0.65, 0.65, 1.55, 1.1, 2.3],
              font_size=10)
    # Work-through box for B
    add_rect(s, 0.7, 4.25, 12.3, 1.7, fill=ACCENT_SFT, line=ACCENT, line_w=0.75)
    add_text(s, "Worked Example — Activity B (UI/UX Design)",
             0.9, 4.3, 6, 0.4, font="Georgia", size=13, bold=True, color=ACCENT)
    calc = ("O = 5 days, M = 10 days, P = 18 days\n"
            "TE = (5 + 4×10 + 18) / 6 = 63 / 6 = 10.5 days\n"
            "σ  = (18 − 5) / 6 = 13 / 6 ≈ 2.17 days\n"
            "Interpretation: ~68% probability of completing within 10.5 ± 2.17 days.")
    add_text(s, calc, 0.9, 4.75, 11.9, 1.1,
             font="Calibri", size=11, color=INK)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 12 — CPM FORWARD + BACKWARD PASS
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = content_slide(prs,
        title="CPM — Forward & Backward Pass",
        eyebrow="Timeline & Milestones",
        notes="The Critical Path Method gives us the longest sequence of dependent activities "
              "and the float available on non-critical ones. Forward pass: ES = max(EF of "
              "predecessors), EF = ES + duration. Backward pass: LF = min(LS of successors), "
              "LS = LF − duration. Slack = LS − ES. Critical activities have zero slack. "
              "Activity F has 8 days of float — it can slip by a week before affecting "
              "the project end date.",
        slide_num=sn, total=TOTAL)
    cpm_data = [
        ["Act.", "Desc.",         "Dur.", "Predecessors",  "ES", "EF", "LS", "LF", "Slack", "CP?"],
        ["A",    "Requirements",   "5 d",  "—",             "0",  "5",  "0",  "5",  "0",     "★"],
        ["B",    "UI/UX Design",  "10 d", "A",             "5",  "15", "5",  "15", "0",     "★"],
        ["C",    "DB Schema",     "4 d",  "A",             "5",  "9",  "14", "18", "9",     ""],
        ["D",    "API Design",    "5 d",  "A",             "5",  "10", "13", "18", "8",     ""],
        ["E",    "Frontend Dev",  "15 d", "B",             "15", "30", "15", "30", "0",     "★"],
        ["F",    "Backend Dev",   "12 d", "C, D",          "10", "22", "18", "30", "8",     ""],
        ["G",    "Integration",   "5 d",  "E, F",          "30", "35", "30", "35", "0",     "★"],
        ["H",    "Testing",       "8 d",  "G",             "35", "43", "35", "43", "0",     "★"],
    ]
    row_fills = [ACCENT_SFT]*2 + [None]*5 + [ACCENT_SFT]*1  # CP rows in accent
    def cp_row_fill(r):
        if r in [0, 1, 4, 6, 7]: return ACCENT_SFT
        return None
    add_table(s, cpm_data, 0.7, 1.55, 12.3, 3.5,
              col_widths=[0.45, 1.5, 0.5, 0.9, 0.45, 0.45, 0.45, 0.45, 0.6, 0.45],
              font_size=10,
              row_fills=[cp_row_fill(r) for r in range(1, len(cpm_data))])

    # ════════════════════════════════════════════════════════════════
    # SLIDE 13 — CRITICAL PATH BIG STATEMENT
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = content_slide(prs,
        title="The Critical Path",
        eyebrow="Timeline & Milestones",
        accent_one=True,
        notes="The critical path is the sequence of activities with zero float — any delay "
              "on them directly delays the project. It runs A → B → E → G → H, totalling "
              "43 working days or approximately 9 calendar weeks. Activity F — Backend "
              "Development — has 8 days of float, giving us a scheduling buffer. The project "
              "duration is dominated by the Frontend Development activity at 15 days.",
        slide_num=sn, total=TOTAL)
    add_rect(s, 0.7, 1.8, 12.3, 0.003, line=RULE, line_w=0.75)
    # Path nodes
    path = [("A", "Requirements\n5 days"), ("B", "UI/UX Design\n10 days"),
            ("E", "Frontend Dev\n15 days"), ("G", "Integration\n5 days"),
            ("H", "Testing\n8 days")]
    for i, (node, lbl) in enumerate(path):
        x = 0.7 + i * 2.5
        add_rect(s, x, 2.0, 2.1, 1.4, fill=ACCENT, line=ACCENT, line_w=1)
        add_text(s, node, x, 2.05, 2.1, 0.7,
                 font="Georgia", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, lbl, x+0.05, 2.75, 2.0, 0.6,
                 font="Calibri", size=11, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(path)-1:
            add_text(s, "→", x+2.1, 2.4, 0.4, 0.6,
                     font="Calibri", size=28, color=ACCENT)
    # Total
    add_rect(s, 0.7, 3.6, 12.3, 0.003, line=RULE, line_w=0.75)
    add_text(s, "A → B → E → G → H  =  5 + 10 + 15 + 5 + 8  =  43 working days",
             0.7, 3.7, 12.3, 0.6, font="Georgia", size=22, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER)
    add_text(s, "≈ 9 calendar weeks  ·  Activity F (Backend Dev) has 8 days of float",
             0.7, 4.35, 12.3, 0.4, font="Calibri", size=14, italic=True, color=INK_SOFT,
             align=PP_ALIGN.CENTER)
    # Float table
    float_data = [
        ["Activity", "Slack", "Interpretation"],
        ["F — Backend Dev",  "8 days", "Can slip up to 8 days without delaying project"],
        ["D — API Design",   "8 days", "Can slip up to 8 days"],
        ["C — DB Schema",    "9 days", "Can slip up to 9 days"],
    ]
    add_table(s, float_data, 0.7, 5.0, 8.0, 1.5,
              col_widths=[2.0, 1.0, 5.0], font_size=10)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 14 — § IV DIVIDER
    # ════════════════════════════════════════════════════════════════
    section_divider(prs, 14, "IV", "Resource Management",
                    "CLO 2, 8, 9 · Ch. 5, 11, 12")
    prs.slides[-1].notes_slide.notes_text_frame.text = (
        "Section IV covers how we allocate people, budget, and materials to the project. "
        "The RACI matrix defines who does what across all 10 major work packages. The "
        "bottom-up budget rolls up from Level 3 hours to a total BAC of 33.6 million VND."
    )

    # ════════════════════════════════════════════════════════════════
    # SLIDE 15 — RACI MATRIX
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = content_slide(prs,
        title="RACI Matrix — Responsibility Assignment",
        eyebrow="Resource Management",
        notes="The RACI matrix answers: who is Accountable, Responsible, Consulted, and "
              "Informed for each work package. Key verification rules from Chapter 12: "
              "exactly one Accountable per row, at least one Responsible per row, and no "
              "team member is purely Informed. Each of our four members carries a mix of "
              "A and R roles — no one is just a passive recipient of information.",
        slide_num=sn, total=TOTAL)
    raci_data = [
        ["Work Package",        "WBS",    "Đức\n(PM)",  "Bảo\n(Dev)",  "Thịnh\n(UI/UX)", "Hiếu\n(QA)"],
        ["Project Planning",     "1.1.1",  "A",          "C",            "I",              "I"],
        ["Requirements",        "1.2",    "A",          "R",            "R",              "C"],
        ["UI/UX Design",        "1.3.1",  "I",          "C",            "A/R",            "I"],
        ["DB Schema",           "1.3.2",  "A",          "R",            "I",              "C"],
        ["API Design",          "1.3.3",  "A",          "R",            "I",              "I"],
        ["Frontend Dev",        "1.4.1",  "C",          "C",            "A/R",            "I"],
        ["Backend Dev",         "1.4.2",  "I",          "A/R",          "I",              "C"],
        ["Unit Testing",        "1.5.1",  "I",          "C",            "I",              "A/R"],
        ["Documentation",       "1.6",    "C",          "C",            "C",              "A/R"],
        ["Presentation",        "1.7",    "A",          "R",            "R",              "R"],
    ]
    add_table(s, raci_data, 0.7, 1.55, 12.3, 4.2,
              col_widths=[2.1, 0.55, 1.55, 1.55, 1.7, 1.55],
              font_size=10)
    add_text(s, "R = Responsible  ·  A = Accountable  ·  C = Consulted  ·  I = Informed",
             0.7, 5.9, 12, 0.3, font="Calibri", size=10, italic=True, color=MUTED)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 16 — BOTTOM-UP BUDGET
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = content_slide(prs,
        title="Bottom-Up Budget Rollup",
        eyebrow="Resource Management",
        notes="The budget is built bottom-up from the WBS. We estimate hours per Level 3 "
              "work package, multiply by the blended labor rate of 50,000 VND per hour, "
              "then add overhead, G&A, contingency, and management reserve. The labor "
              "subtotal of 24 million VND carries a 20% overhead and 10% G&A. The "
              "contingency reserve of 2.4 million is the primary risk buffer, topped up "
              "by a management reserve of 680,000 drawn from the G&A buffer.",
        slide_num=sn, total=TOTAL, accent_one=True)
    budget_data = [
        ["WBS",     "Work Package",             "Hours", "Cost (VND)"],
        ["1.1",     "Project Management",        "60",    "3,000,000"],
        ["1.2",     "Requirements & Analysis",   "50",    "2,500,000"],
        ["1.3",     "System Design",             "80",    "4,000,000"],
        ["1.4",     "Development",              "180",    "9,000,000"],
        ["1.5",     "Testing & QA",              "60",    "3,000,000"],
        ["1.6",     "Documentation",             "30",    "1,500,000"],
        ["1.7",     "Presentation & Submission",  "20",    "1,000,000"],
        ["LABOR",   "Subtotal (480 h × 50K)",   "480",   "24,000,000"],
        ["",        "+ Overhead 20%",            "",      "4,800,000"],
        ["",        "+ G&A 10% (reallocated)",   "",      "1,720,000"],
        ["",        "+ Contingency 10%",          "",      "2,400,000"],
        ["",        "+ Management Reserve",      "",      "680,000"],
        ["BAC",     "Budget at Completion",      "",      "33,600,000 VND"],
    ]
    add_table(s, budget_data, 0.7, 1.55, 8.5, 5.3,
              col_widths=[0.6, 3.4, 0.7, 2.0], font_size=10,
              row_fills=[None]*7 + [BG_ELEV]*1 + [None]*4 + [ACCENT_SFT])
    # Big BAC callout
    add_rect(s, 9.5, 1.55, 3.5, 2.2, fill=ACCENT, line=ACCENT, line_w=1)
    add_text(s, "BAC", 9.6, 1.65, 3.3, 0.5,
             font="Georgia", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "33,600,000", 9.6, 2.15, 3.3, 0.8,
             font="Georgia", size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "VND", 9.6, 2.95, 3.3, 0.4,
             font="Calibri", size=14, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "480 hours × 50,000/h\n+ overhead + contingency",
             9.6, 3.4, 3.3, 0.35,
             font="Calibri", size=9, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 17 — PV CURVE (time-phased budget)
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = content_slide(prs,
        title="Time-Phased Budget — Planned Value Curve",
        eyebrow="Resource Management",
        notes="The PV curve shows cumulative planned spend across the 8-month project timeline. "
              "It's front-loaded — heavy spending in months 4 through 7 during the development "
              "phase, lighter in the planning and final submission months. By month 8, the "
              "cumulative PV reaches the full BAC of 33.6 million VND. This curve was locked "
              "at Week 1 as the cost baseline for EVM tracking.",
        slide_num=sn, total=TOTAL)

    # Build a bar chart for PV using shapes
    pv_data = [
        ("M1",  8,   2_688_000),
        ("M2",  15,  5_040_000),
        ("M3",  25,  8_400_000),
        ("M4",  40, 13_440_000),
        ("M5",  55, 18_480_000),
        ("M6",  70, 23_520_000),
        ("M7",  85, 28_560_000),
        ("M8", 100, 33_600_000),
    ]
    chart_l = 0.7; chart_t = 1.55
    chart_w = 9.5; chart_h = 4.8
    axis_x = chart_l + 0.5; axis_y = chart_t + chart_h - 0.5
    bar_area_w = chart_w - 0.6; bar_area_h = chart_h - 0.7
    max_pct = 100; max_vnd = 33_600_000
    n = len(pv_data)

    # Axes
    add_rect(s, axis_x, axis_y, bar_area_w, 0.003, line=INK, line_w=0.75)  # X axis
    add_rect(s, axis_x, axis_y, 0.003, -bar_area_h, line=INK, line_w=0.75)  # Y axis

    # Y-axis labels (VND)
    for pct in [0, 25, 50, 75, 100]:
        y = axis_y - (pct/100) * bar_area_h
        add_text(s, f"{pct}%", axis_x - 0.45, y - 0.12, 0.42, 0.24,
                 font="Calibri", size=9, color=MUTED, align=PP_ALIGN.RIGHT)
        add_rect(s, axis_x, y, bar_area_w, 0.003, line=RULE, line_w=0.3)
    add_text(s, "VND", axis_x - 0.45, axis_y - bar_area_h - 0.05, 0.42, 0.24,
             font="Calibri", size=9, color=MUTED, align=PP_ALIGN.RIGHT)

    bar_w = (bar_area_w - (n-1)*0.15) / n
    for i, (month, pct, vnd) in enumerate(pv_data):
        bar_h = (pct / max_pct) * bar_area_h
        x = axis_x + i * (bar_w + 0.15)
        y = axis_y - bar_h
        add_rect(s, x, y, bar_w, bar_h,
                 fill=ACCENT if pct == 100 else INK_SOFT, line=None)
        # Cumulative label
        add_text(s, f"{vnd//1000:,}K", x, y - 0.32, bar_w, 0.3,
                 font="Calibri", size=8, color=INK, align=PP_ALIGN.CENTER)
        add_text(s, month, x, axis_y + 0.02, bar_w, 0.24,
                 font="Calibri", size=9, color=INK, align=PP_ALIGN.CENTER)

    # Data table on right
    tbl_data = [
        ["Month", "PV (VND)",  "%"],
        ["M1",   "2,688,000", "8%"],
        ["M2",   "5,040,000", "15%"],
        ["M3",   "8,400,000", "25%"],
        ["M4",   "13,440,000","40%"],
        ["M5",   "18,480,000","55%"],
        ["M6",   "23,520,000","70%"],
        ["M7",   "28,560,000","85%"],
        ["M8",   "33,600,000","100%"],
    ]
    add_table(s, tbl_data, 10.4, 1.55, 2.6, 3.0,
              col_widths=[0.6, 1.15, 0.85], font_size=9)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 18 — § V DIVIDER
    # ════════════════════════════════════════════════════════════════
    section_divider(prs, 18, "V", "Risk Management",
                    "CLO 3, 6 · Ch. 7")
    prs.slides[-1].notes_slide.notes_text_frame.text = (
        "Section V documents our risk management process — identification, assessment, "
        "EMV calculation, and mitigation planning for six key risks. The total risk EMV "
        "of 3.08 million VND drives the contingency reserve. We'll also look at the "
        "3×3 risk heatmap and reconcile EMV against our budgeted contingency."
    )

    # ════════════════════════════════════════════════════════════════
    # SLIDE 19 — RISK REGISTER
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = content_slide(prs,
        title="Risk Register — 6 Risks with EMV",
        eyebrow="Risk Management",
        notes="We identified six risks through team brainstorming and analogy with prior "
              "student projects. Each is scored on likelihood (0–1), with cost and time "
              "impacts. EMV = Likelihood × Impact gives us the expected monetary value. "
              "R2 — scope creep — has the highest EMV at 1.5 million VND, driving our "
              "formal change control process. The total EMV cost of 3.08 million VND and "
              "time EMV of 13.35 days are the key inputs to our contingency planning.",
        slide_num=sn, total=TOTAL)
    risk_data = [
        ["ID", "Risk Description",                    "Cat.",   "L",   "Cost\nImpact",  "Time\nImpact", "EMV Cost",    "EMV Time",  "Strategy",               "Owner"],
        ["R1", "Member unavailable (illness/conflict)", "People","0.30","1,000,000",    "5 d",          "300,000",     "1.50 d",    "Contingency / Accept",    "Đức"],
        ["R2", "Scope creep — new requirements added",  "Scope", "0.50","3,000,000",    "10 d",         "1,500,000",   "5.00 d",    "Avoid + Reduce",         "Đức"],
        ["R3", "Stack unfamiliarity (React 19, TS)",   "Tech",  "0.40","500,000",      "7 d",          "200,000",     "2.80 d",    "Reduce",                 "Bảo"],
        ["R4", "Data loss / version control failure",   "Tech",  "0.15","2,000,000",    "3 d",          "300,000",     "0.45 d",    "Reduce",                 "Thịnh"],
        ["R5", "FE/BE integration failure at W7",       "Tech",  "0.35","2,000,000",    "8 d",          "700,000",     "2.80 d",    "Reduce",                 "Hiếu"],
        ["R6", "Meeting coordination issues",           "People","0.40","200,000",      "2 d",          "80,000",      "0.80 d",    "Accept",                 "Đức"],
        ["",   "TOTAL EMV",                             "",      "",    "",              "",             "3,080,000 VND","13.35 d",   "",                        ""],
    ]
    add_table(s, risk_data, 0.7, 1.55, 12.3, 3.2,
              col_widths=[0.35, 2.3, 0.6, 0.4, 0.85, 0.65, 0.85, 0.65, 1.4, 0.65],
              font_size=9,
              row_fills=[None]*5 + [BG_ELEV]*1 + [ACCENT_SFT])

    # ════════════════════════════════════════════════════════════════
    # SLIDE 20 — RISK HEATMAP
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = content_slide(prs,
        title="Risk Heatmap — Likelihood × Impact",
        eyebrow="Risk Management",
        notes="The heatmap plots each risk on a 3×3 grid of Likelihood versus Impact. "
              "High-likelihood or high-impact risks land in the red zone at top-right — "
              "R2, R3, and R5 require active mitigation. Medium risks R1 and R6 are "
              "monitored. R4 is low-likelihood and low-impact, so we accept it with "
              "preventive measures in place. The heatmap helps prioritise where to focus "
              "contingency effort.",
        slide_num=sn, total=TOTAL)
    # 3×3 grid cells: (likelihood, impact, color, label)
    zone_fill = {
        (3,3): RGBColor(0xA6,0x5D,0x3F), (3,2): RGBColor(0xC4,0x9A,0x7A), (3,1): RGBColor(0xE8,0xD5,0xBC),
        (2,3): RGBColor(0xC4,0x9A,0x7A), (2,2): RGBColor(0xD9,0xCD,0xB8), (2,1): RGBColor(0xF0,0xE9,0xDD),
        (1,3): RGBColor(0xE8,0xD5,0xBC), (1,2): RGBColor(0xF0,0xE9,0xDD), (1,1): RGBColor(0xF5,0xEF,0xE6),
    }
    cell_w = 2.3; cell_h = 1.2
    grid_l = 1.5; grid_t = 1.55
    # Y axis label
    add_text(s, "LIKELIHOOD →", 0.5, 3.5, 0.9, 0.3,
             font="Calibri", size=9, bold=True, color=MUTED,
             align=PP_ALIGN.CENTER)
    # X axis label
    add_text(s, "IMPACT →", 3.2, 6.55, 2.0, 0.3,
             font="Calibri", size=9, bold=True, color=MUTED,
             align=PP_ALIGN.CENTER)
    # Draw grid
    for row in range(1, 4):  # likelihood: 1=Low, 2=Med, 3=High (top to bottom)
        for col in range(1, 4):  # impact: 1=Low, 2=Med, 3=High (left to right)
            x = grid_l + (col-1)*cell_w
            y = grid_t + (3-row)*cell_h
            fill = zone_fill[(row, col)]
            add_rect(s, x, y, cell_w-0.05, cell_h-0.05,
                     fill=fill, line=RULE, line_w=0.5)
    # Risk markers: (likelihood, impact, label, color)
    risks = [(3,3,"R2",ACCENT),(3,2,"R3",ACCENT),(2,3,"R5",ACCENT),
             (2,2,"R1",INK_SOFT),(2,2,"R6",INK_SOFT),(1,1,"R4",MUTED)]
    for (l,i,lbl,col) in risks:
        cx = grid_l + (i-1)*cell_w + cell_w/2 - 0.2
        cy = grid_t + (3-l)*cell_h + cell_h/2 - 0.2
        # Circle marker
        circ = s.shapes.add_shape(9, cx*914400*SW/914400, cy*914400*SH/914400,
                                   Inches(0.38), Inches(0.38))
        circ.fill.solid(); circ.fill.fore_color.rgb = col
        circ.line.color.rgb = BG; circ.line.width = Pt(1)
        add_text(s, lbl, cx, cy, 0.38, 0.38,
                 font="Calibri", size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
    # Zone labels
    add_text(s, "HIGH RISK", grid_l+1.5, grid_t+2.45, 2.5, 0.3,
             font="Calibri", size=9, bold=True, color=WARN, align=PP_ALIGN.CENTER)
    add_text(s, "MEDIUM", grid_l+1.5, grid_t+1.25, 2.5, 0.3,
             font="Calibri", size=9, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, "LOW RISK", grid_l+1.5, grid_t+0.05, 2.5, 0.3,
             font="Calibri", size=9, color=OK, align=PP_ALIGN.CENTER)
    # Legend
    leg_data = [
        ["Risk", "L", "Impact", "Zone"],
        ["R1 — Member unavailable",  "0.30","5 d",   "Medium"],
        ["R2 — Scope creep",         "0.50","10 d",  "High"],
        ["R3 — Stack unfamiliarity",  "0.40","7 d",   "High"],
        ["R4 — Data loss",           "0.15","3 d",   "Low"],
        ["R5 — Integration failure",  "0.35","8 d",   "High"],
        ["R6 — Meeting issues",       "0.40","2 d",   "Medium"],
    ]
    add_table(s, leg_data, 8.5, 1.55, 4.5, 2.8,
              col_widths=[2.0, 0.5, 0.6, 1.4], font_size=9)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 21 — EMV RECONCILIATION
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = content_slide(prs,
        title="EMV Reconciliation",
        eyebrow="Risk Management",
        notes="Our contingency reserve of 2.4 million VND alone doesn't fully cover the "
              "total risk EMV of 3.08 million — there's a gap of 680,000 VND. We address "
              "this by reallocating from the G&A buffer, which was originally budgeted at "
              "2.4 million VND. Reducing G&A to 1.72 million closes the gap while "
              "maintaining full EMV coverage. This decision was made at the Week 3 risk "
              "review and documented in meeting minutes. The total project budget "
              "remains 33.6 million VND — the reserves are reallocated, not added.",
        slide_num=sn, total=TOTAL, accent_one=True)
    recon_data = [
        ["Item",                                      "Amount (VND)", "Notes"],
        ["Total Risk EMV (cost)",                      "3,080,000",    "From risk register Table 14"],
        ["Contingency reserve (10% of labor)",         "2,400,000",    "From budget baseline Table 10"],
        ["Gap (EMV − Contingency)",                    "680,000",      "EMV exceeds contingency"],
        ["Management reserve (from G&A buffer)",       "+ 680,000",    "Drawn from G&A; approved by PM"],
        ["Adjusted G&A (reduced from 2,400,000)",       "1,720,000",    "G&A reduced by 680,000"],
        ["TOTAL RISK COVERAGE",                        "3,080,000 VND","Full EMV covered ✓"],
        ["Budget at Completion (unchanged)",            "33,600,000",   "Total budget unchanged"],
    ]
    add_table(s, recon_data, 0.7, 1.55, 8.0, 3.8,
              col_widths=[3.5, 1.8, 2.7], font_size=11,
              row_fills=[None]*3+[RGBColor(0xF0,0xE9,0xDD)]*1+[None]*2+[ACCENT_SFT])
    # Big callout
    add_rect(s, 9.0, 1.55, 4.0, 2.0, fill=ACCENT, line=ACCENT, line_w=1)
    add_text(s, "EMV = Contingency + Mgmt Reserve",
             9.1, 1.65, 3.8, 0.45, font="Calibri", size=11, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, "3,080,000", 9.1, 2.1, 3.8, 0.7,
             font="Georgia", size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "VND", 9.1, 2.8, 3.8, 0.3,
             font="Calibri", size=14, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "= 2,400,000 + 680,000",
             9.1, 3.1, 3.8, 0.3,
             font="Calibri", size=11, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 22 — § VI DIVIDER
    # ════════════════════════════════════════════════════════════════
    section_divider(prs, 22, "VI", "Communication & Collaboration",
                    "CLO 4, 5, 8, 9, 12 · Ch. 9, 11, 12")
    prs.slides[-1].notes_slide.notes_text_frame.text = (
        "Section VI covers how the team communicates and resolves conflict. We use a "
        "structured 6-step conflict resolution process for interpersonal issues, and "
        "a formal change control process for scope changes. We'll also look at the "
        "stakeholder power/interest grid and our communication channel matrix."
    )

    # ════════════════════════════════════════════════════════════════
    # SLIDE 23 — CONFLICT RESOLUTION + STAKEHOLDER
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = content_slide(prs,
        title="Conflict Resolution & Stakeholder Strategy",
        eyebrow="Communication",
        notes="The 6-step conflict resolution process ensures disagreements are handled "
              "constructively. Key principles: focus on the issue not the person, "
              "document agreements immediately, and escalate only when necessary. "
              "The stakeholder power/interest grid maps five stakeholder groups — "
              "instructors and the university are high power/high interest and require "
              "active engagement. Peer groups and end users need different strategies.",
        slide_num=sn, total=TOTAL)
    # Left: 6-step process
    add_text(s, "6-Step Conflict Resolution", 0.7, 1.55, 6.0, 0.4,
             font="Georgia", size=16, bold=True, color=INK)
    steps = [
        ("1", "Identify root causes",
         "PM calls private session; each party states position in one sentence"),
        ("2", "Gather facts",
         "Review task logs, commit history, WBS — no hearsay"),
        ("3", "Facilitate open discussion",
         "Moderated by PM; time-boxed to 15 min per issue"),
        ("4", "Focus on issues, not personalities",
         "PM enforces ground rules; redirect from 'you' to 'the decision'"),
        ("5", "Seek win-win resolutions",
         "Explore trade-offs; document rationale"),
        ("6", "Document agreements",
         "Added to meeting minutes within 24 h; all parties confirm via Teams"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        y = 2.05 + i * 0.78
        add_rect(s, 0.7, y, 0.4, 0.65, fill=ACCENT, line=None)
        add_text(s, num, 0.7, y, 0.4, 0.65,
                 font="Georgia", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, title, 1.2, y, 2.0, 0.3,
                 font="Calibri", size=11, bold=True, color=INK)
        add_text(s, desc, 1.2, y+0.28, 5.3, 0.4,
                 font="Calibri", size=10, color=INK_SOFT)

    # Right: Power/Interest grid
    add_text(s, "Stakeholder Power / Interest Grid", 7.5, 1.55, 5.5, 0.4,
             font="Georgia", size=16, bold=True, color=INK)
    grid_data = [
        ["",                    "Low Interest",    "High Interest"],
        ["High\nPower",         "Monitor\n(Peer Groups)",   "Keep Satisfied\n(Instructor, Univ.)"],
        ["Low\nPower",          "Keep Informed\n(Bystanders)", "Keep Informed\n(End Users)"],
    ]
    stake_t = s.shapes.add_table(3, 3,
        Inches(7.5), Inches(2.0), Inches(5.4), Inches(3.5)).table
    for r, row in enumerate(grid_data):
        for c, val in enumerate(row):
            cell = stake_t.cell(r, c)
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = val
            run.font.name = "Calibri"
            run.font.size = Pt(10)
            run.font.color.rgb = WHITE if r == 0 else INK
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
            elif r == 1:
                fills = [RGBColor(0xF0,0xE9,0xDD), OK, ACCENT_SFT]
                cell.fill.solid(); cell.fill.fore_color.rgb = fills[c]
            elif r == 2:
                fills = [RGBColor(0xF0,0xE9,0xDD), RGBColor(0xF5,0xEF,0xE6), RGBColor(0xE8,0xD5,0xBC)]
                cell.fill.solid(); cell.fill.fore_color.rgb = fills[c]

    # ════════════════════════════════════════════════════════════════
    # SLIDE 24 — § VII DIVIDER
    # ════════════════════════════════════════════════════════════════
    section_divider(prs, 24, "VII", "Monitoring & EVM",
                    "CLO 7 · Ch. 10")
    prs.slides[-1].notes_slide.notes_text_frame.text = (
        "Section VII is the project control heart — Earned Value Management. At Week 6 "
        "we took a mid-project snapshot. The results show the project is over budget "
        "(CPI = 0.75) and behind schedule (SPI = 0.45). We'll look at the EVM metrics, "
        "the corrective action plan, and the issue log."
    )

    # ════════════════════════════════════════════════════════════════
    # SLIDE 25 — EVM WEEK 6
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = content_slide(prs,
        title="EVM Status — Week 6 Snapshot",
        eyebrow="Monitoring & EVM",
        notes="At the Week 6 milestone checkpoint, we're tracking three metrics. Earned "
              "Value of 7.5 million VND comes from completed milestones: Charter (500K), "
              "Requirements (2.5M), Design Review (4M), and partial prototype (500K). "
              "Actual Cost is 10 million VND from 200 logged hours at 50,000 VND per hour. "
              "Planned Value at Week 6 should be 50% of BAC = 16.8 million VND. The "
              "variance analysis reveals we're both over budget and behind schedule.",
        slide_num=sn, total=TOTAL, accent_one=True)
    # 5 KPI cards
    kpis = [
        ("0.75",   "Cost Performance\nIndex (CPI)",  "Over budget"),
        ("0.45",   "Schedule Perf.\nIndex (SPI)",     "Behind schedule"),
        ("44.8M",  "Est. at Completion\n(EAC)",       "Forecast: 33.6M → 44.8M"),
        ("−2.5M",  "Cost Variance\n(CV)",             "Over budget by 2.5M"),
        ("−11.2M", "Variance at\nCompletion (VAC)",  "Forecast overrun"),
    ]
    for i, (val, lbl, sub) in enumerate(kpis):
        x = 0.7 + i * 2.5
        acc = i == 0
        kpi_card(s, x, 1.55, 2.35, 1.7, lbl, val, sub, accent=acc)
    # EVM table
    evm_data = [
        ["Metric", "Formula",                    "Workings",                           "Value"],
        ["BAC",    "Budget at Completion",       "From budget baseline",              "33,600,000 VND"],
        ["PV",     "Planned Value at W6",        "50% of BAC",                         "16,800,000 VND"],
        ["EV",     "Earned Value",                "M1+M2+M3+partial = 500K+2.5M+4M+0.5M","7,500,000 VND"],
        ["AC",     "Actual Cost",                 "200 h × 50,000 VND/h",              "10,000,000 VND"],
        ["CV",     "Cost Variance",               "EV − AC = 7.5M − 10M",              "−2,500,000 VND"],
        ["SV",     "Schedule Variance",            "EV − PV = 7.5M − 16.8M",           "−9,300,000 VND"],
        ["CPI",    "Cost Perf. Index",             "EV / AC = 7.5M / 10M",              "0.75"],
        ["SPI",    "Schedule Perf. Index",         "EV / PV = 7.5M / 16.8M",           "≈ 0.45"],
        ["EAC",    "Estimate at Completion",       "BAC / CPI = 33.6M / 0.75",         "44,800,000 VND"],
        ["ETC",    "Estimate to Complete",         "EAC − AC = 44.8M − 10M",           "34,800,000 VND"],
        ["VAC",    "Variance at Completion",       "BAC − EAC = 33.6M − 44.8M",        "−11,200,000 VND"],
    ]
    add_table(s, evm_data, 0.7, 3.4, 12.3, 3.6,
              col_widths=[0.7, 2.0, 4.1, 2.0], font_size=9,
              row_fills=[None]*4+[WARN]*2+[None]*5+[WARN])

    # ════════════════════════════════════════════════════════════════
    # SLIDE 26 — ISSUE LOG + CHANGE CONTROL
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = content_slide(prs,
        title="Issue Log & Change Control",
        eyebrow="Monitoring & EVM",
        notes="The issue log tracks problems as they arise. ISS-001 was closed — stakeholder "
              "interviews ran 3 hours over. ISS-002 was closed — we downgraded react-router-dom "
              "from v7 to v6 as a workaround. ISS-003 is open — the Week 10 integration "
              "milestone is at risk. The change control table shows CR-001, which approved "
              "replacing Supabase with local state to simplify the prototype for the demo. "
              "This had a positive schedule impact and minor cost saving.",
        slide_num=sn, total=TOTAL)
    # Issue log
    add_text(s, "Issue Log", 0.7, 1.55, 5.8, 0.4,
             font="Georgia", size=15, bold=True, color=INK)
    issue_data = [
        ["ID",     "Week", "Description",                            "Owner","Status"],
        ["ISS-001","W3",   "Stakeholder interviews took +3 h",       "Đức",  "Closed"],
        ["ISS-002","W5",   "react-router-dom v7 routing broken",     "Bảo",  "Closed"],
        ["ISS-003","W6",   "M4 Integration milestone at risk",       "Đức",  "Open"],
    ]
    add_table(s, issue_data, 0.7, 2.0, 5.8, 1.8,
              col_widths=[0.65, 0.5, 3.0, 0.65, 0.65], font_size=9,
              row_fills=[None, RGBColor(0xE8,0xF0,0xE4)]*2+[RGBColor(0xF5,0xE9,0xDD)])
    # Change control
    add_text(s, "Change Request CR-001", 6.8, 1.55, 6.2, 0.4,
             font="Georgia", size=15, bold=True, color=INK)
    cr_data = [
        ["Field",           "Value"],
        ["Date",            "Week 6"],
        ["Requester",       "Bảo"],
        ["Description",     "Replace Supabase (cloud) with local state + mock data"],
        ["Impact — Time",   "+1 day float recovered; no schedule impact"],
        ["Impact — Cost",   "−200,000 VND (eliminates Supabase setup)"],
        ["Impact — Scope",  "Removes cloud sync; already in Out-of-Scope"],
        ["Decision",        "APPROVED"],
        ["Status",          "IMPLEMENTED"],
    ]
    add_table(s, cr_data, 6.8, 2.0, 6.2, 3.6,
              col_widths=[1.6, 4.6], font_size=9,
              row_fills=[None]*7+[OK]+[OK])

    # ════════════════════════════════════════════════════════════════
    # SLIDE 27 — § VIII DIVIDER
    # ════════════════════════════════════════════════════════════════
    section_divider(prs, 27, "VIII", "Prototype & Lessons",
                    "CLO 7, 8, 9 · Ch. 13")
    prs.slides[-1].notes_slide.notes_text_frame.text = (
        "Section VIII closes the loop — we show the live demo plan for three prototype "
        "features, then reflect on five key lessons learned tagged to their relevant CLOs. "
        "We'll also confirm full CLO coverage across all nine learning outcomes."
    )

    # ════════════════════════════════════════════════════════════════
    # SLIDE 28 — PROTOTYPE DEMO PLAN
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = content_slide(prs,
        title="Prototype Demonstration Plan",
        eyebrow="Prototype & Lessons",
        notes="We demonstrate three core features in under 3 minutes. Feature 1 is Task "
              "CRUD with assignment — showing how to create a task, assign it to a team "
              "member, set priority and deadline. Feature 2 is deadline management with "
              "overdue highlighting — the system automatically flags past-deadline tasks "
              "in red. Feature 3 is the progress dashboard — project-level and task-level "
              "metrics aggregated in real time. The demo environment is localhost:5173.",
        slide_num=sn, total=TOTAL)
    features = [
        ("Feature 1", "Task CRUD + Assignment", "90 s",
         "Navigate to /tasks\nCreate task: 'Design landing page',\nassignee = 'Bảo', deadline = June 20,\npriority = High\n→ Verify task appears with assignee badge",
         "WBS 1.4.1", "Bảo + Thịnh"),
        ("Feature 2", "Deadline Management\n+ Overdue Highlight", "60 s",
         "Show task with past deadline\n(yesterday: June 15)\n→ Red overdue badge displayed\n→ Filter by 'Overdue'",
         "WBS 1.3.1", "Thịnh"),
        ("Feature 3", "Progress Dashboard", "60 s",
         "Navigate to /dashboard\n→ Project cards with counts\n→ Task progress bar (e.g., 60%)\n→ Real-time metric aggregation",
         "WBS 1.4.1, 1.4.2", "Bảo"),
    ]
    cols = [0.7, 4.7, 8.7]
    for i, (feat, title, time, script, wbs, owner) in enumerate(features):
        x = cols[i]
        add_rect(s, x, 1.55, 3.8, 0.45, fill=ACCENT, line=None)
        add_text(s, f"{feat}  ·  {time}", x, 1.55, 3.8, 0.45,
                 font="Calibri", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(s, x, 2.05, 3.8, 4.5, fill=BG_ELEV, line=RULE, line_w=0.75)
        add_text(s, title, x+0.1, 2.1, 3.6, 0.65,
                 font="Georgia", size=14, bold=True, color=ACCENT)
        add_text(s, "Script:", x+0.1, 2.8, 3.6, 0.3,
                 font="Calibri", size=10, bold=True, color=MUTED)
        add_text(s, script, x+0.1, 3.1, 3.6, 2.0,
                 font="Calibri", size=10, color=INK)
        add_text(s, f"WBS: {wbs}", x+0.1, 5.5, 3.6, 0.3,
                 font="Calibri", size=10, color=MUTED)
        add_text(s, f"Owner: {owner}", x+0.1, 5.85, 3.6, 0.3,
                 font="Calibri", size=10, color=MUTED)
    add_text(s, "Demo URL: http://localhost:5173  ·  Run: npm install && npm run dev",
             0.7, 7.0, 12.3, 0.3,
             font="Calibri", size=11, italic=True, color=MUTED)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 29 — LESSONS LEARNED
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = content_slide(prs,
        title="Lessons Learned",
        eyebrow="Prototype & Lessons",
        notes="Five concrete lessons from the project. LL-1: scope creep consumed 4 hours "
              "before CR-001 formalised the scope lock — we now require a change request "
              "form from day one. LL-2: PERT estimates were too optimistic — Activity B "
              "took 14 days vs. 10.5 estimated. Future projects should use a 1.2× "
              "multiplier for design-phase estimates. LL-3: react-router v7 broke our "
              "build — we now pin versions in package.json. LL-4: EVM was reactive; "
              "we'll implement weekly mini-checks going forward. LL-5: parallel dev "
              "without a frozen API contract caused 2 days of rework at integration.",
        slide_num=sn, total=TOTAL)
    lessons = [
        ("LL-1", "Scope creep at W6",
         "CR-001 formalised scope lock; 4 h recovered",
         "CLO 3, 9"),
        ("LL-2", "PERT estimates too optimistic",
         "Activity B took 14 days; TE=10.5. Use 1.2× multiplier.",
         "CLO 2, 5"),
        ("LL-3", "react-router v7 incompatibility",
         "Downgraded to v6; 3 rollback commits. Pin versions.",
         "CLO 3, 5"),
        ("LL-4", "EVM was reactive, not proactive",
         "Week 6 VAC = −11.2M. Implement weekly EVM mini-check.",
         "CLO 6, 7"),
        ("LL-5", "Parallel dev without frozen API contract",
         "2 days rework at W7 integration. Freeze API at W5.",
         "CLO 2, 4"),
    ]
    ll_data = [["ID", "Lesson", "Evidence / Action", "CLO Tags"]] + \
               [[a,b,c,d] for a,b,c,d in lessons]
    add_table(s, ll_data, 0.7, 1.55, 12.3, 3.0,
              col_widths=[0.55, 2.2, 6.8, 0.85], font_size=10,
              row_fills=[None]*5)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 30 — CLO COVERAGE MATRIX
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = content_slide(prs,
        title="CLO Coverage Matrix — All 9 CLOs",
        eyebrow="Prototype & Lessons",
        notes="Every one of the nine INS3044 Course Learning Outcomes is demonstrably "
              "covered in this report. CLO 1 through CLO 9 map to specific sections and "
              "evidence artefacts. The coverage check confirms full CLO achievement across "
              "scope, time, cost, risk, quality, communication, and project delivery.",
        slide_num=sn, total=TOTAL)
    clo_data = [
        ["CLO",  "Description",                                       "Section",  "Evidence"],
        ["CLO 1","APPLY PM principles",                               "§ 1",      "SMART objectives; full PM cycle"],
        ["CLO 2","IMPLEMENT scope / HR / time / cost",                "§ 1–4",    "WBS, CPM, RACI, bottom-up budget"],
        ["CLO 3","IMPLEMENT integration / risk / quality",             "§ 5, § 8", "Risk EMV; quality metrics; CR-001"],
        ["CLO 4","IMPLEMENT systems & decisions",                      "§ 6, § 8", "Decision process; live demo"],
        ["CLO 5","USE tools / processes / techniques",                 "§ 2–3",    "WBS tool, MS Project, PERT/CPM"],
        ["CLO 6","PERFORM risk analysis & contingency",                "§ 5",      "6 risks; EMV; contingency reconciliation"],
        ["CLO 7","PERFORM overall evaluation",                        "§ 7",      "10 EVM metrics; 5 lessons tagged"],
        ["CLO 8","IMPLEMENT overall project plans",                    "All",      "Report is the project plan; ZIP package"],
        ["CLO 9","Autonomy & personal qualities",                       "§ 6, § 8", "4 members with A/R roles; conflict process"],
    ]
    add_table(s, clo_data, 0.7, 1.55, 12.3, 4.2,
              col_widths=[0.65, 3.2, 0.8, 3.2], font_size=10)
    add_rect(s, 0.7, 5.9, 12.3, 0.003, line=RULE, line_w=0.75)
    add_text(s, "Coverage check: All 9 CLOs are demonstrably addressed through specific report sections and evidence. ✓",
             0.7, 6.0, 12.3, 0.4, font="Calibri", size=12, bold=True, color=OK)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 31 — THANK YOU / Q&A
    # ════════════════════════════════════════════════════════════════
    sn = next_num()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    shp = s.shapes.add_shape(1, 0, 0, SW, SH)
    shp.fill.solid(); shp.fill.fore_color.rgb = BG; shp.line.fill.background()
    add_rect(s, 0, 0, 1, 0.005, line=RULE, line_w=1)
    # Large thank-you
    add_text(s, "Thank you", 0.9, 1.2, 11.5, 1.4,
             font="Georgia", size=72, bold=True, italic=True, color=INK)
    add_text(s, "Questions & Discussion",
             0.9, 2.7, 11.5, 0.7,
             font="Georgia", size=32, italic=True, color=INK_SOFT)
    add_rect(s, 0.9, 3.45, 11.533, 0.003, line=RULE, line_w=1)
    add_text(s, "Group 14 · INS3044 IT Project Management",
             0.9, 3.6, 11.5, 0.4, font="Calibri", size=14, color=MUTED)
    # Team members
    for i, (role, name, sid) in enumerate(members):
        x = 0.9 + i * 2.9
        add_rect(s, x, 4.1, 2.7, 1.1, fill=BG_ELEV, line=RULE, line_w=0.75)
        add_text(s, role,  x+0.1, 4.15, 2.5, 0.3, font="Calibri", size=9,  color=MUTED)
        add_text(s, name,  x+0.1, 4.4,  2.5, 0.35, font="Georgia", size=14, bold=True, color=INK)
        add_text(s, sid,   x+0.1, 4.8,  2.5, 0.3, font="Calibri", size=10, color=MUTED)
    add_rect(s, 0.9, 5.5, 11.533, 0.003, line=RULE, line_w=0.75)
    add_text(s, "Nguyễn Long Đức (23070435) · Phạm Hồ Bảo (23070455) · Kiều Bá Thịnh (23070247) · Đỗ Huy Hiếu (23070325)",
             0.9, 5.6, 11.5, 0.4, font="Calibri", size=11, color=INK_SOFT)
    add_text(s, "Faculty of Applied Sciences · Vietnam National University · June 2026",
             0.9, 6.05, 11.5, 0.3, font="Calibri", size=11, color=MUTED)
    s.notes_slide.notes_text_frame.text = (
        "Thank you for your attention. We're happy to take questions on any aspect of "
        "the project — from the PERT calculations to the EVM corrective action plan. "
        "All four team members can answer questions in their respective areas: Đức on "
        "schedule and risk, Bảo on technical implementation, Thịnh on design and UX, "
        "and Hiếu on testing and quality. The full report, prototype, and presentation "
        "materials are available in the submission ZIP."
    )

    # ── Save ─────────────────────────────────────────────────────────────────
    out_path = os.path.join(os.path.dirname(__file__), "Group14_INS3044_Cartesian.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")
    return out_path

if __name__ == "__main__":
    build()
